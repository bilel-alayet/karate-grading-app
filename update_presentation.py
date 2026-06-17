import sys
import subprocess
import os

print("Installing playwright if needed...")
try:
    import playwright
except ImportError:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "playwright"])
    subprocess.check_call([sys.executable, "-m", "playwright", "install", "chromium"])

from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

screenshot_path = "screenshot.png"

print("Taking screenshot of the web app...")
# Take screenshot
with sync_playwright() as p:
    browser = p.chromium.launch(headless=True)
    page = browser.new_page()
    try:
        page.goto("https://karatesport.web.app", timeout=30000)
        # wait for flutter to render
        page.wait_for_timeout(4000) 
        page.screenshot(path=screenshot_path)
        print("Screenshot saved.")
    except Exception as e:
        print("Error taking screenshot:", e)
    finally:
        browser.close()

print("Updating presentation...")
pptx_path = r"c:\Users\USER\Downloads\Application de Grading de Karaté.pptx"
if not os.path.exists(pptx_path):
    print("PPTX not found:", pptx_path)
    sys.exit(1)

prs = Presentation(pptx_path)

# Try to find Démonstration slide
demo_slide = None
for slide in prs.slides:
    if slide.shapes.title and "Démonstration" in slide.shapes.title.text:
        demo_slide = slide
        break

if demo_slide and os.path.exists(screenshot_path):
    # clear existing bullet points in placeholder if any, or just overlay
    if len(demo_slide.placeholders) > 1:
        demo_slide.placeholders[1].text = ""
    demo_slide.shapes.add_picture(screenshot_path, Inches(2), Inches(2), height=Inches(4.5))

# Add Code samples slide 1: Auth and Firestore
slide_layout = prs.slide_layouts[1]
code_slide1 = prs.slides.add_slide(slide_layout)
code_slide1.shapes.title.text = "Exemple de Code : Authentification & Firestore"
tf1 = code_slide1.placeholders[1].text_frame
tf1.clear()
p1 = tf1.paragraphs[0]
p1.text = """// Inscription avec Firebase Authentication
final cred = await FirebaseAuth.instance.createUserWithEmailAndPassword(
  email: _ec.text.trim(),
  password: _pc.text,
);
await cred.user?.updateDisplayName(_nc.text.trim());

// Sauvegarde du profil dans Cloud Firestore
await FirebaseFirestore.instance.collection('Users').doc(cred.user!.uid).set({
  'email': _ec.text.trim(),
  'nom': _nc.text.trim(),
  'role': 'parent',
});"""
p1.font.size = Pt(14)
p1.font.name = "Consolas"

# Add Code samples slide 2: Evaluation Firestore
code_slide2 = prs.slides.add_slide(slide_layout)
code_slide2.shapes.title.text = "Exemple de Code : Sauvegarde d'Évaluation"
tf2 = code_slide2.placeholders[1].text_frame
tf2.clear()
p2 = tf2.paragraphs[0]
p2.text = """// Ajout des notes de karaté dans la base de données
await FirebaseFirestore.instance.collection('inscriptions').add({
  'enfantNom': _nc.text.trim(),
  'age': int.parse(_ac.text.trim()),
  'ceinture': _ceinture, // ex: Blanche, Jaune...
  'parentUid': widget.parentUid,
  'statut': 'En attente',
  'createdAt': FieldValue.serverTimestamp(),
});"""
p2.font.size = Pt(14)
p2.font.name = "Consolas"

prs.save(pptx_path)
print("Successfully updated PPTX:", pptx_path)
